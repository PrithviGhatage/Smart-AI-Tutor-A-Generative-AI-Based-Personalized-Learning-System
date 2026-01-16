import streamlit as st
from openai import OpenAI
from pptx import Presentation

# ================== THEME COLORS ==================
PRIMARY = "#3B82F6"      # main accent (blue)
SECONDARY = "#22D3EE"    # neon-ish accent
BACKGROUND = "#020617"   # main background (dark)
SURFACE = "#020617"      # cards / boxes
TEXT_PRIMARY = "#E5E7EB" # main text color (light grey)
TEXT_MUTED = "#6B7280"   # secondary text


# ================== PAGE CONFIG ==================
st.set_page_config(
    page_title="Smart AI Tutor",
    layout="wide",
)

# ================== GLOBAL STYLE ==================
st.markdown(
    f"""
<style>
/* Replace Streamlit’s default red tab underline with blue */
.stTabs [data-baseweb="tab-highlight"] {{
    background-color: {PRIMARY} !important;  /* blue underline */
    height: 3px !important;
    border-radius: 10px !important;
}}
/* Slider theme (remove red, apply blue theme) */
.stSlider > div > div > div > div {{
    background-color: {PRIMARY} !important;  /* slider knob */
}}

.stSlider > div > div > div {{
    background-color: #1e3a8a !important;    /* slider track */
}}


.stApp {{
    background-color: {BACKGROUND};
    color: {TEXT_PRIMARY};
    font-family: system-ui, -apple-system, BlinkMacSystemFont, "SF Pro Text", "Segoe UI", sans-serif;
}}
/* Slider knob (handle) – force blue instead of red */
.stSlider [role="slider"] {{
    background-color: {PRIMARY} !important;   /* blue circle */
    border: 2px solid #e5e7eb !important;     /* light outline */
}}
/* Slider knob (handle) – color + number text */
.stSlider [role="slider"] {{
    background-color: {PRIMARY} !important;   /* blue knob */
    border: 2px solid #e5e7eb !important;     /* light border */
}}

/* Any text inside the knob (the number) */
.stSlider [role="slider"] * {{
    color: #F9FAFB !important;                /* white number */
}}


/* Layout width & spacing */
.block-container {{
    padding-top: 1.5rem;
    padding-bottom: 1.5rem;
    padding-left: 3rem;
    padding-right: 3rem;
    max-width: 1150px;
}}

/* Headings */
h1 {{
    font-size: 2.4rem;
    font-weight: 900;
    letter-spacing: 0.04em;
    color: #F9FAFB;
    margin-bottom: 0.2rem;
}}

/* Subheadings etc. */
h2, h3, h4, h5, h6 {{
    color: {TEXT_PRIMARY};
    font-weight: 700;
}}

/* Body text */
p, span, label {{
    color: {TEXT_MUTED};
    font-size: 0.95rem;
}}

/* Card container (chat-style) */
section.main > div > div > div {{
    background-color: {SURFACE};
    border-radius: 16px;
    border: 1px solid #111827;
    padding: 18px 22px;
}}

/* Inputs */
.stTextInput > div > div > input,
.stTextArea textarea,
textarea,
input[type="text"] {{
    background-color: #020617 !important;
    color: {TEXT_PRIMARY} !important;
    border-radius: 10px !important;
    border: 1px solid #1f2937 !important;
}}

.stTextInput > div > div > input:focus,
.stTextArea textarea:focus,
textarea:focus,
input[type="text"]:focus {{
    border-color: {SECONDARY} !important;
    box-shadow: 0 0 0 1px {SECONDARY}55 !important;
}}

/* Buttons */
.stButton > button,
div[data-testid="stButton"] > button {{
    background-color: {PRIMARY} !important;     /* blue */
    color: white !important;
    border-radius: 10px !important;
    border: none !important;
    padding: 0.55rem 1.4rem !important;
    font-weight: 600 !important;
    font-size: 0.95rem !important;
    transition: 0.15s ease-in-out;
}}

/* Hover effect */
.stButton > button:hover,
div[data-testid="stButton"] > button:hover {{
    background-color: #1d4ed8 !important;       /* darker blue */
    box-shadow: 0 0 12px #3b82f655 !important;
    color: white !important;
}}

/* BUTTON TEXT FIX — forces white text even if Streamlit wraps it */
.stButton > button p,
div[data-testid="stButton"] > button p {{
    color: white !important;
    font-weight: 600 !important;
}}

/* ================= SIDEBAR BUTTON — RESET CONVERSATION ================= */

.stSidebar .stButton > button,
.stSidebar div[data-testid="stButton"] > button {{
    background-color: {PRIMARY} !important;     /* same blue */
    color: white !important;
    border-radius: 10px !important;
    border: none !important;
    font-weight: 600 !important;
}}

.stSidebar .stButton > button:hover,
.stSidebar div[data-testid="stButton"] > button:hover {{
    background-color: #1d4ed8 !important;
    box-shadow: 0 0 12px #3b82f655 !important;
}}

/* Make all radio buttons use blue accent instead of red */
input[type="radio"] {{
    accent-color: {PRIMARY};
}}

/* Radio buttons – use blue theme instead of red */
.stRadio [role="radio"] > div:first-child {{
    width: 14px;
    height: 14px;
    border-radius: 999px;
    border: 2px solid {PRIMARY} !important;   /* blue outline */
    background-color: transparent !important;
    box-shadow: none !important;
}}

/* Selected radio – filled blue */
.stRadio [role="radio"][aria-checked="true"] > div:first-child {{
    background-color: {PRIMARY} !important;   /* blue fill */
}}
/* FORCE radio buttons to use blue fill and blue border */
.stRadio [role="radio"] {{
    border: none !important;
}}

.stRadio [role="radio"] > div:first-child {{
    width: 16px !important;
    height: 16px !important;
    border-radius: 50% !important;
    border: 2px solid {PRIMARY} !important;    /* blue ring */
    background-color: transparent !important;
}}

/* The INNER DOT for selected radio */
.stRadio [role="radio"][aria-checked="true"] > div:first-child {{
    background-color: {PRIMARY} !important;    /* filled blue */
    border: 2px solid {PRIMARY} !important;
}}
/* Make radio button dot blue */
input[type="radio"] {{
    accent-color: {PRIMARY} !important;
}}
/* Make the start of the slider (selected range) blue */
.stSlider div[data-baseweb="slider"] div[aria-hidden="true"] {{
    background-color: {PRIMARY} !important;
}}


</style>
""",
    unsafe_allow_html=True,
)

# ================== SIMPLE HEADER ==================
st.title("Smart AI Tutor")
st.caption("Final Year Capstone · AI tutor that learns from PPTs and evaluates your answers")

# ================== SESSION STATE ==================
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

if "ppt_text" not in st.session_state:
    st.session_state.ppt_text = None

if "ppt_title" not in st.session_state:
    st.session_state.ppt_title = None

if "eval_questions" not in st.session_state:
    st.session_state.eval_questions = []  # list of questions

if "eval_current_index" not in st.session_state:
    st.session_state.eval_current_index = 0  # which question we are on

if "eval_results" not in st.session_state:
    st.session_state.eval_results = []  # list of dicts: {q, answer, evaluation}


# ================== SIDEBAR SETTINGS ==================
st.sidebar.header("⚙️ Settings")

api_key = st.sidebar.text_input(
    "OpenAI API Key",
    type="password",
    help="Your key is only used locally on your machine.",
)

subject = st.sidebar.selectbox(
    "Subject / Domain",
    [
        "Programming",
        "Data Structures & Algorithms",
        "Databases",
        "Operating Systems",
        "Computer Networks",
        "Mathematics",
        "Machine Learning",
        "Other",
    ],
)

level = st.sidebar.selectbox(
    "Student Level",
    ["School", "College", "Competitive Exams", "Beginner Programmer", "Intermediate Programmer"],
)

tutor_style = st.sidebar.radio(
    "Tutor Style",
    ["Friendly & simple", "Exam focused", "Deep conceptual"],
    index=1,
)

if st.sidebar.button("🔁 Reset Conversation"):
    st.session_state.chat_history = []
    st.success("Conversation reset.")

st.sidebar.markdown("---")
st.sidebar.markdown("**Built with Python · Streamlit · OpenAI**")

# ================== OPENAI HELPER ==================
def get_client(api_key: str) -> OpenAI:
    return OpenAI(api_key=api_key)


def call_chat_model(client: OpenAI, system_prompt: str, history, user_message: str) -> str:
    """
    history: list of {"role": "user"/"assistant", "content": "..."}
    """
    messages = [{"role": "system", "content": system_prompt}]
    messages.extend(history)
    messages.append({"role": "user", "content": user_message})

    response = client.chat.completions.create(
        model="gpt-4o-mini",  # change to gpt-4o if you have access
        messages=messages,
        temperature=0.7,
    )
    return response.choices[0].message.content


# ================== SYSTEM PROMPT BUILDER ==================
def build_system_prompt() -> str:
    style_instructions = {
        "Friendly & simple": (
            "Use very simple language, small steps, and examples from daily life. "
            "Regularly ask short questions to check understanding."
        ),
        "Exam focused": (
            "Focus on definitions, key points, and exam-oriented language. "
            "Structure answers with headings, bullet points, and short summaries."
        ),
        "Deep conceptual": (
            "Give rigorous, detailed explanations with analogies, small proofs, and why/how. "
            "Encourage the student to think and do small derivations."
        ),
    }

    return f"""
You are an AI Tutor for the subject: {subject}.
The student level is: {level}.

Your goals:
1. Explain concepts clearly and step-by-step.
2. Adapt explanations to the given level.
3. Use examples related to the subject.
4. Encourage active learning by asking quick check questions sometimes.
5. Avoid giving full direct solutions immediately for homework/exam-like questions: first guide, then reveal.

Tutor style instructions:
{style_instructions[tutor_style]}

Very IMPORTANT:
- Always keep answers structured with headings and bullet points where useful.
- At the end of each answer, add a short section:
  '✅ Quick Summary' with 3–5 bullet points.
"""
# ================== PPT HELPER ==================
def extract_text_from_pptx(file) -> str:
    """
    Takes a PPTX file-like object and returns all text from slides as one big string.
    """
    prs = Presentation(file)
    texts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)

    return "\n\n".join(texts)

# ================== TABS LAYOUT ==================

tab_chat, tab_explain, tab_quiz, tab_exam, tab_ppt, tab_eval = st.tabs(
    [
        "Tutor Chat",
        "Concept Explainer",
        "Practice Quiz",
        "Exam Style Answer",
        "Learn from PPT",
        "Practice & Evaluation",
    ]
)


# ---------- TAB 1: Tutor Chat ----------
with tab_chat:
    st.subheader("Tutor Chat")

    user_input = st.text_area(
        "Your question:",
        placeholder="Example: Explain dynamic programming in simple words.",
        height=120,
    )

    if st.button("Ask Tutor"):
        if not api_key:
            st.error("Please enter your OpenAI API key in the sidebar.")
        elif not user_input.strip():
            st.warning("Please type a question.")
        else:
            client = get_client(api_key)
            system_prompt = build_system_prompt()

            with st.spinner("Generating response..."):
                reply = call_chat_model(
                    client=client,
                    system_prompt=system_prompt,
                    history=st.session_state.chat_history,
                    user_message=user_input,
                )

            st.session_state.chat_history.append(
                {"role": "user", "content": user_input}
            )
            st.session_state.chat_history.append(
                {"role": "assistant", "content": reply}
            )

    if st.session_state.chat_history:
        st.markdown("---")
        st.markdown("### Conversation History")
        for msg in st.session_state.chat_history:
            if msg["role"] == "user":
                st.markdown(f"**You:** {msg['content']}")
            else:
                st.markdown(f"**Tutor:** {msg['content']}")

# ---------- TAB 2: Concept Explainer ----------
with tab_explain:
    st.subheader("Concept Explainer")

    topic = st.text_input(
        "Enter a topic or concept",
        placeholder="Example: Binary Search Trees, Normalization, Gradient Descent",
    )

    depth = st.slider(
        "Detail Level (1 = simple, 5 = very detailed)",
        min_value=1,
        max_value=5,
        value=3,
    )

    style_choice = st.selectbox(
        "Explanation Style",
        [
            "Simple explanation",
            "College exam explanation",
            "Interview/viva explanation",
        ],
    )

    if st.button("Explain"):
        if not api_key:
            st.error("Please enter your OpenAI API key in the sidebar.")
        elif not topic.strip():
            st.warning("Please enter a topic.")
        else:
            client = get_client(api_key)
            base_prompt = build_system_prompt()

            extra_instruction = f"""
The student wants a focused explanation of the topic: {topic}.
Explanation style: {style_choice}.
Depth level: {depth}.

Structure your answer:
1. Intuition or analogy
2. Formal definition
3. Step-by-step explanation
4. Small example
5. Common mistakes
6. Summary (3-5 key points)
"""

            with st.spinner("Generating explanation..."):
                answer = call_chat_model(
                    client,
                    system_prompt=base_prompt + extra_instruction,
                    history=[],
                    user_message=f"Explain the topic '{topic}' as requested.",
                )

            st.markdown("### Explanation")
            st.write(answer)

# ---------- TAB 3: Practice Quiz ----------
with tab_quiz:
    st.subheader("Practice Quiz")

    quiz_topic = st.text_input(
        "Topic for quiz",
        placeholder="Example: SQL Joins, Time Complexity, Linked Lists",
    )

    num_questions = st.slider(
        "Number of questions",
        min_value=3,
        max_value=15,
        value=5,
    )

    q_type = st.selectbox(
        "Question type",
        ["Mixed (MCQ + short answer)", "Only MCQ", "Only short answer"],
    )

    if st.button("Generate Quiz"):
        if not api_key:
            st.error("Please enter your OpenAI API key in the sidebar.")
        elif not quiz_topic.strip():
            st.warning("Please enter a quiz topic.")
        else:
            client = get_client(api_key)
            base_prompt = build_system_prompt()

            quiz_instruction = f"""
You are in QUIZ mode.

Create {num_questions} questions for the topic: {quiz_topic}
Subject: {subject}
Level: {level}
Question type preference: {q_type}

For each question:
- Number it: Q1, Q2, ...
- If MCQ, give 4 options (A, B, C, D).
- After all questions, provide an 'Answer Key and Explanations' section.

Output format:
Q1) ...
A) ...
B) ...
C) ...
D) ...

Q2) ...

---

Answer Key and Explanations:
Q1) Correct option: X
Short explanation: ...
"""

            with st.spinner("Generating quiz..."):
                quiz_text = call_chat_model(
                    client,
                    system_prompt=base_prompt + quiz_instruction,
                    history=[],
                    user_message="Generate the quiz as described.",
                )

            st.markdown("### Generated Quiz")
            st.write(quiz_text)

# ---------- TAB 4: Exam Style Answer ----------
with tab_exam:
    st.subheader("Exam Style Answer Generator")

    exam_question = st.text_area(
        "Enter an exam-style question",
        placeholder="Example: Explain normalization in DBMS. Discuss 1NF, 2NF, 3NF with examples.",
        height=140,
    )

    marks = st.selectbox(
        "Marks for this question",
        [3, 5, 8, 10, 15],
        index=2,
    )

    if st.button("Generate Exam Answer"):
        if not api_key:
            st.error("Please enter your OpenAI API key in the sidebar.")
        elif not exam_question.strip():
            st.warning("Please enter a question.")
        else:
            client = get_client(api_key)
            base_prompt = build_system_prompt()

            exam_instruction = f"""
You are writing an exam answer for the question below.

Marks: {marks}
Subject: {subject}
Level: {level}

Rules:
- Write like a high-scoring university exam answer.
- Use clear headings and bullet points.
- Include definitions, diagrams (described in text), and examples where helpful.
- Keep the length appropriate for the marks.
- End with a short summary section.
"""

            with st.spinner("Generating exam-style answer..."):
                exam_answer = call_chat_model(
                    client,
                    system_prompt=base_prompt + exam_instruction,
                    history=[],
                    user_message=exam_question,
                )

            st.markdown("### Model Exam Answer")
            st.write(exam_answer)
# ---------- TAB 5: Learn from PPT ----------
with tab_ppt:
    st.subheader("Learn Directly from PPT Slides")

    st.markdown(
        "Upload a PPTX file from your class. "
        "The tutor will read the slides, summarise them, and answer questions using only that content."
    )

    uploaded_ppt = st.file_uploader(
        "Upload PPTX file",
        type=["pptx"],
        help="Download the PPT from email and upload it here.",
    )

    col_p1, col_p2 = st.columns(2)
    with col_p1:
        summarize_clicked = st.button("Summarise Slides")
    with col_p2:
        learn_clicked = st.button("Ask Questions from PPT")

    if uploaded_ppt is not None:
        st.session_state.ppt_title = uploaded_ppt.name
        st.session_state.ppt_text = extract_text_from_pptx(uploaded_ppt)
        st.success(f"PPT loaded: {uploaded_ppt.name}")

        with st.expander("Preview extracted text (for checking)"):
            if st.session_state.ppt_text:
                preview_text = (
                    st.session_state.ppt_text[:2000] + "..."
                    if len(st.session_state.ppt_text) > 2000
                    else st.session_state.ppt_text
                )
                st.text(preview_text)

    # Summarise PPT
    if summarize_clicked:
        if not api_key:
            st.error("Please enter your OpenAI API key in the sidebar.")
        elif not st.session_state.ppt_text:
            st.warning("Please upload a PPT first.")
        else:
            client = get_client(api_key)
            base_prompt = build_system_prompt()
            context_text = st.session_state.ppt_text[:8000]

            instruction = f"""
You are given text extracted from a PowerPoint presentation.

File name: {st.session_state.ppt_title}

Using only this content, create:
1. A high-level summary (5–10 bullet points).
2. A list of key concepts, formulas, or definitions.
3. 5 possible viva questions based on the slides.
4. A short summary at the end.

Here is the PPT content:
\"\"\"{context_text}\"\"\"
"""

            with st.spinner("Summarising slides..."):
                summary = call_chat_model(
                    client,
                    system_prompt=base_prompt + instruction,
                    history=[],
                    user_message="Summarise this PPT for the student.",
                )

            st.markdown("### PPT Summary and Key Points")
            st.write(summary)

    st.markdown("---")

    # Q&A based only on PPT
    ppt_question = st.text_area(
        "Ask a question based on this PPT",
        placeholder="Example: Explain the slide on normalization, or 'What are the main topics in this PPT?'",
        height=120,
    )

    if learn_clicked:
        if not api_key:
            st.error("Please enter your OpenAI API key in the sidebar.")
        elif not st.session_state.ppt_text:
            st.warning("Please upload a PPT first.")
        elif not ppt_question.strip():
            st.warning("Please type a question.")
        else:
            client = get_client(api_key)
            base_prompt = build_system_prompt()
            context_text = st.session_state.ppt_text[:8000]

            qa_instruction = f"""
You are an AI tutor that must answer using only the content from this PPT.

File name: {st.session_state.ppt_title}

Here is the PPT content:
\"\"\"{context_text}\"\"\"

Rules:
- If the answer is clearly in the PPT content, explain it clearly.
- If the answer is not clearly mentioned, say:
  "This specific point is not clearly mentioned in your slides."
"""

            with st.spinner("Answering using PPT content..."):
                answer = call_chat_model(
                    client,
                    system_prompt=base_prompt + qa_instruction,
                    history=[],
                    user_message=ppt_question,
                )

            st.markdown("### Answer Based on PPT")
            st.write(answer)
# ---------- TAB 6: Practice and Evaluation ----------
with tab_eval:
    st.subheader("Practice and Evaluation")

    st.markdown(
        "In this mode, the system generates practice questions for a topic. "
        "You answer them one by one, and the AI evaluates your answer and gives a score and feedback."
    )

    # Step 1: Generate questions
    with st.expander("Step 1: Generate practice questions", expanded=True):
        eval_topic = st.text_input(
            "Topic for practice",
            placeholder="Example: Normalization in DBMS, Binary Search, Transactions in DBMS",
        )
        eval_num_q = st.slider(
            "Number of questions",
            min_value=3,
            max_value=10,
            value=5,
        )

        if st.button("Generate Practice Questions"):
            if not api_key:
                st.error("Please enter your OpenAI API key in the sidebar.")
            elif not eval_topic.strip():
                st.warning("Please enter a topic.")
            else:
                client = get_client(api_key)
                base_prompt = build_system_prompt()

                gen_instruction = f"""
You are an AI tutor generating practice questions for active recall.

Topic: {eval_topic}
Subject: {subject}
Level: {level}
Number of questions: {eval_num_q}

Rules:
- Only short-answer questions (no MCQs).
- Each question should be answerable in 3–6 sentences.
- Make difficulty gradually increase.
- Output format, one per line:
Q1) ...
Q2) ...
Q3) ...
"""

                with st.spinner("Generating questions..."):
                    questions_text = call_chat_model(
                        client,
                        system_prompt=base_prompt + gen_instruction,
                        history=[],
                        user_message="Generate only the questions.",
                    )

                lines = questions_text.split("\n")
                questions = [ln.strip() for ln in lines if ln.strip().startswith("Q")]
                if not questions and questions_text.strip():
                    questions = [questions_text.strip()]

                st.session_state.eval_questions = questions
                st.session_state.eval_current_index = 0
                st.session_state.eval_results = []

                st.success("Practice questions generated.")

    # Show overview of questions
    if st.session_state.eval_questions:
        st.markdown("### Generated Questions")
        for q in st.session_state.eval_questions:
            st.markdown(f"- {q}")

    st.markdown("---")

    # Step 2: Answer and evaluation
    if not st.session_state.eval_questions:
        st.info("Generate questions above to start a practice session.")
    else:
        idx = st.session_state.eval_current_index
        total = len(st.session_state.eval_questions)

        if idx < total:
            current_q = st.session_state.eval_questions[idx]
            st.markdown(f"### Question {idx + 1} of {total}")
            st.markdown(f"**{current_q}**")

            student_answer = st.text_area(
                "Your answer",
                height=160,
                placeholder="Write your answer here in your own words.",
                key=f"answer_{idx}",
            )

            if st.button("Submit Answer for Evaluation"):
                if not api_key:
                    st.error("Please enter your OpenAI API key in the sidebar.")
                elif not student_answer.strip():
                    st.warning("Please write an answer.")
                else:
                    client = get_client(api_key)
                    base_prompt = build_system_prompt()

                    eval_instruction = f"""
You are a strict but helpful university examiner.

Evaluate the student's answer to the question below.

Question:
{current_q}

Student's answer:
\"\"\"{student_answer}\"\"\"

Subject: {subject}
Level: {level}

Tasks:
1. Give a score from 0 to 10.
2. Explain briefly why you gave that score.
3. Provide a model answer in 4–8 bullet points.
4. Give 2–3 suggestions for improvement.

Format:

Score: X/10

Feedback:
- ...

Model Answer:
- ...

Suggestions:
- ...
"""

                    with st.spinner("Evaluating answer..."):
                        evaluation = call_chat_model(
                            client,
                            system_prompt=base_prompt + eval_instruction,
                            history=[],
                            user_message="Evaluate the student's answer as requested.",
                        )

                    st.session_state.eval_results.append(
                        {
                            "question": current_q,
                            "answer": student_answer,
                            "evaluation": evaluation,
                        }
                    )
                    st.session_state.eval_current_index += 1

                    st.markdown("### Evaluation")
                    st.write(evaluation)

        # Session summary when done
        if (
            st.session_state.eval_current_index >= len(st.session_state.eval_questions)
            and st.session_state.eval_results
        ):
            st.markdown("## Practice Session Summary")

            for i, res in enumerate(st.session_state.eval_results, start=1):
                st.markdown(f"### Question {i}")
                st.markdown(f"**{res['question']}**")
                with st.expander("Your answer"):
                    st.write(res["answer"])
                with st.expander("Evaluation and model answer"):
                    st.write(res["evaluation"])

            if st.button("Start New Practice Session"):
                st.session_state.eval_questions = []
                st.session_state.eval_current_index = 0
                st.session_state.eval_results = []
                st.success("Session reset. You can generate new questions now.")

